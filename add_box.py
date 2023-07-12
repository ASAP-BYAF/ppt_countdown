import pptx
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from color_code_to_RGB import color_code_to_RGB_func


def add_box_func(sld, fontsize=Pt(200), text=None, text_RGB=RGBColor(0, 0, 0), \
                 bg_RGB=None, lm=Pt(0), tm=Pt(0), h=Pt(1000), w=Pt(1000) ):

    # textbox のサイズを指定(height, width は適当でもよさそう)
    left = lm
    top = tm
    height = h
    width = w

    # テキストボックスを追加
    shp = sld.shapes.add_textbox(left, top, width, height)

    # 背景色の設定
    if bg_RGB:
        shp.fill.solid()
        shp.fill.fore_color.rgb = bg_RGB

    # 文字の設定
    if text:
        tfrm = shp.text_frame
        pg = tfrm.paragraphs[0]
        pg.text = text
        pg.font.size = fontsize
        pg.alignment = PP_ALIGN.CENTER
        pg.font.color.rgb = text_RGB
        tfrm.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ### fig1
    # # 丸角四角形の挿入
    # fig1 = sld.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    # # 塗りつぶしの色を青色に指定
    # # fig1.fill.solid()
    # # fig1.fill.fore_color.rgb = RGBColor(0, 0, 255)
    # # # 枠線の色を青色に指定
    # # fig1.line.color.rgb = RGBColor(0, 0, 255)
    # # 文字の挿入
    # pg = fig1.text_frame.paragraphs[0]
    # pg.text = "SAMPLE"
    # pg.font.size = Pt(50)
    return sld
